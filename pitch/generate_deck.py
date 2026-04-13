"""
Sessio Pitch Deck Generator — 6 slides, 3-minute pitch
Clean dark theme, minimal text, big numbers
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Theme colors ---
BG_DARK = RGBColor(0x0F, 0x17, 0x29)       # Deep navy
BG_CARD = RGBColor(0x1A, 0x24, 0x3B)       # Slightly lighter card
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
ACCENT = RGBColor(0x4E, 0xA8, 0xF5)        # Blue accent
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)  # Green for positive
ACCENT_RED = RGBColor(0xF4, 0x72, 0x72)    # Red for negative/loss
ACCENT_YELLOW = RGBColor(0xFB, 0xBF, 0x24) # Yellow for highlights
DIM = RGBColor(0x64, 0x70, 0x87)           # Dimmed text


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, default_size=18,
                  default_color=WHITE, font_name="Calibri"):
    """Add a text box with multiple formatted lines.
    Each line is a tuple: (text, font_size, color, bold, alignment)
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold, align = line_data, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else False
            align = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(6)

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# CREATE PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ============================================================
# SLIDE 1: PROBLEM
# ============================================================

slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BG_DARK)

# Small "SESSIO" brand top-left
add_text(slide1, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

# Hook quote — the hotel room analogy
add_multiline(slide1, 0.6, 1.0, 12, 1.5, [
    ('"A coaching slot is like a hotel room \u2014', 32, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ('if it\'s empty tonight, the revenue is gone forever."', 32, WHITE, True, PP_ALIGN.LEFT),
])

# Two problem cards side by side
# Card 1: Admin time
card1 = add_rounded_rect(slide1, 0.6, 3.2, 5.5, 2.5, BG_CARD)
add_multiline(slide1, 1.0, 3.4, 4.7, 2.2, [
    ("\u23F0  ADMIN OVERHEAD", 14, ACCENT, True),
    ("", 8, BG_CARD),
    ("1 day / week", 40, ACCENT_RED, True),
    ("", 8, BG_CARD),
    ("spent on WhatsApp confirmations", 18, LIGHT_GRAY),
    ("and Excel attendance tracking", 18, LIGHT_GRAY),
])

# Card 2: Empty slots
card2 = add_rounded_rect(slide1, 7.0, 3.2, 5.5, 2.5, BG_CARD)
add_multiline(slide1, 7.4, 3.4, 4.7, 2.2, [
    ("\U0001F4B8  EMPTY SLOTS", 14, ACCENT, True),
    ("", 8, BG_CARD),
    ("\u20AC1\u20132K / month", 40, ACCENT_RED, True),
    ("", 8, BG_CARD),
    ("lost to last-minute cancellations", 18, LIGHT_GRAY),
    ("that no one fills", 18, LIGHT_GRAY),
])

# Bottom line
add_text(slide1, 0.6, 6.2, 12, 0.6,
         "No tool on the market solves both.",
         22, ACCENT_YELLOW, True, PP_ALIGN.LEFT)

add_notes(slide1, """[SLIDE 1 — PROBLEM, 0:00-0:30]

"A coaching slot is like a hotel room — if it's empty tonight, the revenue is gone forever.

A small sports school owner spends up to one full day per week managing confirmations on WhatsApp and tracking attendance in Excel. And even after all that work, they still lose one to two thousand euros a month to empty slots from late cancellations.

The admin eats their time. The cancellations eat their revenue. And no tool on the market solves both."
""")


# ============================================================
# SLIDE 2: SOLUTION
# ============================================================

slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BG_DARK)

add_text(slide2, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

# Headline
add_text(slide2, 0.6, 0.8, 12, 0.8,
         "One tool that saves time AND fills revenue gaps.",
         34, WHITE, True)

# One-liner
add_text(slide2, 0.6, 1.7, 11, 0.8,
         "Sessio replaces WhatsApp scheduling for small sports schools and fills "
         "their empty slots by offering last-minute cancellations as discounted "
         "bookings to nearby athletes.",
         18, LIGHT_GRAY)

# 3-step flow
# Step 1
step1 = add_rounded_rect(slide2, 0.6, 3.2, 3.6, 2.8, BG_CARD)
add_multiline(slide2, 0.9, 3.4, 3.0, 2.5, [
    ("\u2776", 28, ACCENT, True),
    ("MANAGE", 22, WHITE, True),
    ("", 6, BG_CARD),
    ("Groups, confirmations,", 16, LIGHT_GRAY),
    ("attendance \u2014 all in one app.", 16, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("Replaces WhatsApp + Excel", 14, ACCENT, False),
])

# Arrow
add_text(slide2, 4.35, 4.0, 0.6, 0.8, "\u279C", 36, DIM, False, PP_ALIGN.CENTER)

# Step 2
step2 = add_rounded_rect(slide2, 4.9, 3.2, 3.6, 2.8, BG_CARD)
add_multiline(slide2, 5.2, 3.4, 3.0, 2.5, [
    ("\u2777", 28, ACCENT_YELLOW, True),
    ("CANCEL", 22, WHITE, True),
    ("", 6, BG_CARD),
    ("Someone can't make it?", 16, LIGHT_GRAY),
    ("Spot opens instantly.", 16, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("No more lost revenue", 14, ACCENT_YELLOW, False),
])

# Arrow
add_text(slide2, 8.65, 4.0, 0.6, 0.8, "\u279C", 36, DIM, False, PP_ALIGN.CENTER)

# Step 3
step3 = add_rounded_rect(slide2, 9.2, 3.2, 3.6, 2.8, BG_CARD)
add_multiline(slide2, 9.5, 3.4, 3.0, 2.5, [
    ("\u2778", 28, ACCENT_GREEN, True),
    ("FILL", 22, WHITE, True),
    ("", 6, BG_CARD),
    ("Offered at a discount to", 16, LIGHT_GRAY),
    ("nearby athletes.", 16, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("Empty slot = new client", 14, ACCENT_GREEN, False),
])

# Bottom punchline
add_text(slide2, 0.6, 6.4, 12, 0.6,
         "Your wasted capacity IS your marketing.",
         22, ACCENT_YELLOW, True)

add_notes(slide2, """[SLIDE 2 — SOLUTION, 0:30-1:00]

"Sessio does two things. First, it replaces WhatsApp and Excel — confirmations, attendance, cancellations, all in one app.

Second — and this is the key — when a slot goes empty, Sessio offers it as a discounted booking to athletes nearby. The coach earns something instead of nothing. The athlete discovers a new coach at low risk. The empty slot becomes a marketing channel.

One tool that saves time AND fills revenue gaps."
""")


# ============================================================
# SLIDE 3: WHY NOW + WEDGE
# ============================================================

slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, BG_DARK)

add_text(slide3, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

add_text(slide3, 0.6, 0.8, 12, 0.8,
         "Why now? The thesis is funded. Europe is empty.",
         34, WHITE, True)

# Left side: Why now bullets
add_rounded_rect(slide3, 0.6, 2.0, 5.8, 3.0, BG_CARD)
add_multiline(slide3, 1.0, 2.2, 5.0, 2.6, [
    ("WHY NOW", 14, ACCENT, True),
    ("", 6, BG_CARD),
    ("Coaching demand tripled since 2018", 20, WHITE, True),
    ("", 4, BG_CARD),
    ("TeachMe.To: $10M raised (Sam Altman) \u2192 US only", 17, LIGHT_GRAY),
    ("", 4, BG_CARD),
    ("Booksy: Polish-founded, $1.3B GMV", 17, LIGHT_GRAY),
    ("\u2192 proved the playbook for beauty", 17, LIGHT_GRAY),
    ("", 4, BG_CARD),
    ("Europe \u2192 wide open", 20, ACCENT_GREEN, True),
])

# Right side: Wedge diagram
add_rounded_rect(slide3, 7.0, 2.0, 5.8, 3.0, BG_CARD)
add_multiline(slide3, 7.4, 2.2, 5.0, 2.6, [
    ("THE WEDGE", 14, ACCENT, True),
    ("", 6, BG_CARD),
    ("To fill slots in real time, you need BOTH:", 17, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("ActiveNow = tool, no marketplace", 18, ACCENT_RED),
    ("", 3, BG_CARD),
    ("TeachMe.To = marketplace, no tool", 18, ACCENT_RED),
    ("", 3, BG_CARD),
    ("Sessio = both", 24, ACCENT_GREEN, True),
    ("", 3, BG_CARD),
    ("Neither can add the other half", 16, DIM),
    ("without rebuilding from scratch.", 16, DIM),
])

# Bottom
add_text(slide3, 0.6, 5.5, 12, 0.8,
         "Someone will build this in Europe in the next 18 months. We already started.",
         20, ACCENT_YELLOW, True)

add_notes(slide3, """[SLIDE 3 — WHY NOW + WEDGE, 1:00-1:30]

"Coaching demand tripled in six years, but schools still run on WhatsApp. TeachMe.To just raised ten million dollars from Sam Altman for a coaching marketplace — but they're US only, and after one session coaches take clients off-platform. Europe is wide open. Booksy proved this playbook works from Poland at over a billion in transactions.

And our structural advantage: to fill empty slots in real time, you need to be both the scheduling tool and the marketplace. ActiveNow has the tool but no marketplace. TeachMe.To has the marketplace but no tool. Neither can build the other half without starting over."
""")


# ============================================================
# SLIDE 4: BUSINESS MODEL + MARKET
# ============================================================

slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, BG_DARK)

add_text(slide4, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

add_text(slide4, 0.6, 0.8, 12, 0.8,
         "The Booksy playbook. Proven at $1.3B GMV.",
         34, WHITE, True)

# Booksy comparison table
# Header row
add_rounded_rect(slide4, 0.6, 2.0, 12.0, 0.6, ACCENT)
add_text(slide4, 0.8, 2.05, 3.5, 0.5, "", 16, WHITE, True)
add_text(slide4, 4.2, 2.05, 3.5, 0.5, "Booksy (beauty)", 16, BG_DARK, True, PP_ALIGN.CENTER)
add_text(slide4, 8.2, 2.05, 4.0, 0.5, "Sessio (sports coaching)", 16, BG_DARK, True, PP_ALIGN.CENTER)

# Row 1: SaaS
add_rounded_rect(slide4, 0.6, 2.7, 12.0, 0.65, BG_CARD)
add_text(slide4, 0.8, 2.75, 3.3, 0.5, "SaaS subscription", 16, ACCENT, True)
add_text(slide4, 4.2, 2.75, 3.5, 0.5, "$30/month", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide4, 8.2, 2.75, 4.0, 0.5, "School pays for tool", 16, WHITE, False, PP_ALIGN.CENTER)

# Row 2: Commission
add_rounded_rect(slide4, 0.6, 3.45, 12.0, 0.65, RGBColor(0x14, 0x1E, 0x33))
add_text(slide4, 0.8, 3.5, 3.3, 0.5, "New client commission", 16, ACCENT, True)
add_text(slide4, 4.2, 3.5, 3.5, 0.5, "30% first visit", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide4, 8.2, 3.5, 4.0, 0.5, "30% happy hour booking", 16, WHITE, False, PP_ALIGN.CENTER)

# Row 3: After first visit
add_rounded_rect(slide4, 0.6, 4.2, 12.0, 0.65, BG_CARD)
add_text(slide4, 0.8, 4.25, 3.3, 0.5, "After first visit", 16, ACCENT, True)
add_text(slide4, 4.2, 4.25, 3.5, 0.5, "0% forever", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text(slide4, 8.2, 4.25, 4.0, 0.5, "0% forever", 16, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# Booksy result
add_text(slide4, 0.6, 5.2, 12, 0.5,
         "Booksy result: $66M revenue (2024), $1.3B GMV, profitable. From Poland.",
         18, ACCENT_YELLOW, True, PP_ALIGN.LEFT)

# Market numbers at bottom
add_rounded_rect(slide4, 0.6, 5.9, 3.6, 1.1, BG_CARD)
add_multiline(slide4, 0.9, 5.95, 3.0, 1.0, [
    ("$15.4B", 28, ACCENT, True),
    ("US coaching market", 14, LIGHT_GRAY),
])

add_rounded_rect(slide4, 4.5, 5.9, 3.6, 1.1, BG_CARD)
add_multiline(slide4, 4.8, 5.95, 3.0, 1.0, [
    ("24.5% CAGR", 28, ACCENT, True),
    ("platform layer growth", 14, LIGHT_GRAY),
])

add_rounded_rect(slide4, 8.4, 5.9, 4.2, 1.1, BG_CARD)
add_multiline(slide4, 8.7, 5.95, 3.6, 1.0, [
    ("2,100+ schools", 28, ACCENT, True),
    ("in Poland already pay for tools", 14, LIGHT_GRAY),
])

add_notes(slide4, """[SLIDE 4 — BUSINESS MODEL + MARKET, 1:30-2:00]

"We follow the Booksy model. SaaS subscription for the scheduling tool — schools pay because it saves them a day per week. Commission only on new clients the platform brings — thirty percent on the first visit, zero after that. Nothing to route around.

Booksy built this into sixty-six million in revenue from Poland. The coaching market is fifteen billion dollars in the US, the platform layer is growing twenty-five percent per year, and over two thousand Polish schools already pay for scheduling tools. The market is there. The tools are just bad."
""")


# ============================================================
# SLIDE 5: TRACTION + TEAM
# ============================================================

slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, BG_DARK)

add_text(slide5, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

add_text(slide5, 0.6, 0.8, 12, 0.8,
         "Built. Validated. Ready to scale.",
         34, WHITE, True)

# Left: Traction
add_rounded_rect(slide5, 0.6, 2.0, 5.8, 4.8, BG_CARD)
add_multiline(slide5, 1.0, 2.2, 5.0, 4.4, [
    ("TRACTION", 14, ACCENT, True),
    ("", 10, BG_CARD),
    ("\u2705  MVP built and live", 20, WHITE, True),
    ("", 6, BG_CARD),
    ("\u2705  10+ discovery calls", 20, WHITE, True),
    ("     across 3 countries (PL, US, UA)", 16, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("\u2705  Pain confirmed at every school at scale", 20, WHITE, True),
    ("", 6, BG_CARD),
    ("\u2705  2 pilot-ready schools", 20, WHITE, True),
    ("     400+ clients, 10-16 coaches each", 16, LIGHT_GRAY),
    ("", 6, BG_CARD),
    ("\u2705  WTP confirmed: \u20AC50-100/month", 20, WHITE, True),
])

# Right: Team
add_rounded_rect(slide5, 7.0, 2.0, 5.8, 4.8, BG_CARD)
add_multiline(slide5, 7.4, 2.2, 5.0, 4.4, [
    ("FOUNDER", 14, ACCENT, True),
    ("", 10, BG_CARD),
    ("Miroslaw Bartkowiak", 24, WHITE, True),
    ("", 8, BG_CARD),
    ("Amazon  \u2014  BI Engineer", 18, LIGHT_GRAY),
    ("data systems at $75M/year scale", 16, DIM),
    ("", 4, BG_CARD),
    ("BCG X  \u2014  2 years", 18, LIGHT_GRAY),
    ("client delivery & product", 16, DIM),
    ("", 8, BG_CARD),
    ("Built MVP solo with AI", 18, ACCENT_GREEN, True),
    ("", 4, BG_CARD),
    ("\U0001F1FA\U0001F1E6 \U0001F1F5\U0001F1F1 \U0001F1EC\U0001F1E7 \U0001F1EB\U0001F1F7  4 languages \u2014 built for CEE", 17, LIGHT_GRAY),
    ("", 4, BG_CARD),
    ("Mom runs a swimming school", 17, ACCENT_YELLOW),
    ("\u2192 saw the problem firsthand", 16, DIM),
])

add_notes(slide5, """[SLIDE 5 — TRACTION + TEAM, 2:00-2:30]

"The MVP is built and live. Over ten discovery calls across three countries — pain confirmed at every school at scale. Two tennis schools with four hundred-plus clients ready to pilot.

Before Sessio, I was at Amazon building data systems at seventy-five million dollar scale, and two years at BCG X. I built the MVP solo with AI. I speak four languages and can sell across Central and Eastern Europe from day one. My mom runs a swimming school — that's where I first saw this problem."
""")


# ============================================================
# SLIDE 6: THE ASK
# ============================================================

slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, BG_DARK)

add_text(slide6, 0.6, 0.4, 2, 0.4, "SESSIO", 14, DIM, True)

# Big number
add_text(slide6, 0.6, 1.0, 12, 1.2,
         "\u20AC300\u2013500K pre-seed",
         48, WHITE, True, PP_ALIGN.CENTER)

add_text(slide6, 0.6, 2.2, 12, 0.6,
         "Warsaw is the starting point. Europe is the market.",
         22, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Use of funds — 3 boxes
add_rounded_rect(slide6, 0.6, 3.3, 3.6, 1.8, BG_CARD)
add_multiline(slide6, 0.9, 3.45, 3.0, 1.5, [
    ("60%", 36, ACCENT, True),
    ("PRODUCT", 16, WHITE, True),
    ("Hire 1-2 developers", 14, LIGHT_GRAY),
    ("Ship faster", 14, LIGHT_GRAY),
])

add_rounded_rect(slide6, 4.8, 3.3, 3.6, 1.8, BG_CARD)
add_multiline(slide6, 5.1, 3.45, 3.0, 1.5, [
    ("25%", 36, ACCENT, True),
    ("GO-TO-MARKET", 16, WHITE, True),
    ("Onboard 50 schools", 14, LIGHT_GRAY),
    ("in Warsaw", 14, LIGHT_GRAY),
])

add_rounded_rect(slide6, 9.0, 3.3, 3.6, 1.8, BG_CARD)
add_multiline(slide6, 9.3, 3.45, 3.0, 1.5, [
    ("15%", 36, ACCENT, True),
    ("OPERATIONS", 16, WHITE, True),
    ("Runway & infra", 14, LIGHT_GRAY),
    ("18-24 months", 14, LIGHT_GRAY),
])

# Milestones
add_rounded_rect(slide6, 0.6, 5.5, 12.0, 1.6, BG_CARD)
add_multiline(slide6, 1.0, 5.6, 11.2, 1.4, [
    ("MILESTONES TO SEED", 14, ACCENT, True),
    ("", 6, BG_CARD),
    ("50 paying schools in Warsaw    \u00B7    Prove daily usage retention    \u00B7    "
     "Validate happy hours fill slots    \u00B7    Seed-ready",
     18, WHITE, False, PP_ALIGN.CENTER),
])

# Thank you
add_text(slide6, 0.6, 7.0, 12, 0.4,
         "Thank you.  \u2014  mbartkovyak@gmail.com",
         16, DIM, False, PP_ALIGN.CENTER)

add_notes(slide6, """[SLIDE 6 — THE ASK, 2:30-3:00]

"We're raising three to five hundred thousand euros pre-seed. Sixty percent goes to product — hiring one or two developers to move faster. Twenty-five percent goes to go-to-market — onboarding fifty paying schools in Warsaw.

Our goal: fifty schools using Sessio daily, proof that discounted bookings fill empty slots, and enough traction to raise a seed. Warsaw is the starting point. Europe is the market. Thank you."
""")


# ============================================================
# SAVE
# ============================================================

output_path = "/Users/myro/Documents/Sessio/pitch/sessio_pitch.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
